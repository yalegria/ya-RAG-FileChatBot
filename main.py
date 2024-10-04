import streamlit as st

# Library for data extraction,analysis, conversion & manipulation of PDFs. 
import fitz  # PyMuPDF

from io import BytesIO

# TODO: Deprecated, migrate to langchain_openai.ChatOpenAI
from langchain_community.chat_models import ChatOpenAI

from langchain.schema.messages import HumanMessage
import pandas as pd
from docx import Document

# Optical character recognition library
import pytesseract

from PIL import Image
from pptx import Presentation
import csv
from zipfile import BadZipFile  # Import BadZipFile from zipfile module
import os

# Initialize LangChain models
chain_gpt_35 = ChatOpenAI(model="gpt-3.5-turbo")

# Initialize session state attributes
if 'history' not in st.session_state:
    st.session_state.history = []
if 'files_uploaded' not in st.session_state:
    st.session_state.files_uploaded = []
if 'file_contents' not in st.session_state:
    st.session_state.file_contents = []  # Store extracted text content from files
if 'file_titles' not in st.session_state:
    st.session_state.file_titles = {}  # Store file titles

# Title
st.set_page_config(
    page_title="RAG File Assistant - Yuri Alegria",
    page_icon=":robot:",  # You can use emoji or the path to an image file
)


# Sidebar for API key input
#with st.sidebar:
#    api_key = st.text_input("Enter your OpenAI API Key", type="password")
#    if api_key:
#        os.environ["OPENAI_API_KEY"] = api_key


# Assign files to variables and titles
def assign_files_to_vars(files):
    file_vars = {}
    for i, file in enumerate(files):
        file_key = f'file{i+1}'
        file_vars[file_key] = file
        st.session_state.file_titles[file_key] = file.name  # Store file title
    return file_vars



# Function to append response to chat history
def append_response_to_history(user_query, response):
    # Append both user query and assistant response to history
    st.session_state.history.append({"user": user_query, "assistant": response.content})

# Function to summarize content if it's too long
def summarize_content(content):
    # If the content length exceeds 2048 characters, summarize it
    if len(content) > 16384:
        response = chain_gpt_35.invoke([HumanMessage(content=f"Summarize the following content:\n\n{content[:16000]}")])
        return response.content
    return content

# Function to process files
def process_files(files):
    data_found = False
    file_vars = assign_files_to_vars(files)
    for file_key, file in file_vars.items():
        file_type = file.type
        if file_type in ['image/png', 'image/jpeg']:
            data_found |= img_reader(file, file_key)
        elif file_type == 'application/pdf':
            data_found |= pdf_reader(file, file_key)
        elif file_type == 'text/plain':
            data_found |= txt_reader(file, file_key)
        elif file_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
            data_found |= doc_reader(file, file_key)
        elif file_type in ['application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
            data_found |= excel_reader(file, file_key)
        elif file_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
            data_found |= ppt_reader(file, file_key)
        elif file_type == 'text/csv':
            data_found |= csv_reader(file, file_key)
    
    if not data_found:
        #st.warning("The provided files do not meet your requirements.")  # Use warning instead of adding to history
        pass

# File Readers with content extraction

# Image reader for image files containing text
def img_reader(file, file_key):
    img = Image.open(file)
    text = pytesseract.image_to_string(img)
    st.session_state.file_contents.append(f"{file_key}: {summarize_content(text)}")

    return True

# PDF reader using fitz library
def pdf_reader(file, file_key):
    try:
        file_stream = BytesIO(file.read())
        if file_stream.getbuffer().nbytes == 0:
            #st.error("The provided PDF file is empty.")  # Use error instead of adding to history
            return False
        
        pdf_document = fitz.open(stream=file_stream, filetype="pdf")
        full_text = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            full_text += page.get_text()
        pdf_document.close()
        
        if full_text.strip() == "":
            st.error("The provided PDF file does not contain readable text.")  # Use error instead of adding to history
            return False
        
        st.session_state.file_contents.append(f"{file_key}: {summarize_content(full_text)}")
        return True

    except fitz.EmptyFileError:
        st.error("The provided PDF file is empty or corrupted.")  # Use error instead of adding to history
        return False

def txt_reader(file, file_key):
    text = file.read().decode("utf-8")
    if text.strip() == "":
        return False
    st.session_state.file_contents.append(f"{file_key}: {summarize_content(text)}")
    return True

def doc_reader(file, file_key):
    try:
        doc = Document(BytesIO(file.read()))
        text = "\n".join([para.text for para in doc.paragraphs])
        if text.strip() == "":
            #st.error("The provided Word document is empty or does not contain readable text.")  # Use error instead of adding to history
            return False
        st.session_state.file_contents.append(f"{file_key}: {summarize_content(text)}")
        return True

    except BadZipFile:
        st.error("The provided file could not be read as a valid Word document (BadZipFile).")  # Use error instead of adding to history
        return False

def excel_reader(file, file_key):
    df = pd.read_excel(file)
    text = df.to_string()
    if text.strip() == "":
        return False
    st.session_state.file_contents.append(f"{file_key}: {summarize_content(text)}")
    return True

def ppt_reader(file, file_key):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    if text.strip() == "":
        return False
    st.session_state.file_contents.append(f"{file_key}: {summarize_content(text)}")
    return True

def csv_reader(file, file_key):
    df = pd.read_csv(file)
    text = df.to_string()
    if text.strip() == "":
        return False
    st.session_state.file_contents.append(f"{file_key}: {summarize_content(text)}")
    return True

# Function to handle user queries and invoke models
def handle_user_query(user_query):
    combined_content = "\n\n".join(st.session_state.file_contents)  # Combine all file contents
    summarized_content = summarize_content(combined_content)  # Summarize combined content
    prompt = f"Based on the following uploaded document contents:\n\n{summarized_content}\n\nUser's question: {user_query}"

    # Transform a single input into an output
    # HumanMessage represents a message with role "user"
    response = chain_gpt_35.invoke([HumanMessage(content=prompt)])


    append_response_to_history(user_query, response)


# Display chat history

#def display_history():
#    for message in st.session_state.history[-5:]:  # Display only the last 5 messages to manage token limit
#        if message['user']:  # Only display messages that are not empty or system messages
#            st.write(f"**User:** {message['user']}")
#            st.write(f"**Assistant:** {message['assistant']}")


def display_history():
    for message in st.session_state.history[-5:]:  # Display only the last 5 messages to manage token limit
        if message['user']:  # Only display messages that are not empty or system messages
            st.markdown(f"""
                <p style="margin-bottom:2px;font-weight:bold;">User:</p>
                <div style='background-color:#d1f0ff; padding:10px; border-radius:10px; margin-bottom:10px;'>
                    <p>{message['user']}</p>
                </div>
            """, unsafe_allow_html=True)
            
            st.markdown(f"""
                <p style="margin-bottom:2px;font-weight:bold;">Assistant:</p>
                <div style='background-color:#f7f7f7; padding:10px; border-radius:10px; margin-bottom:10px;'>
                    <p>{message['assistant']}</p>
                </div>
            """, unsafe_allow_html=True)


# Intro page layout (centered header, input, and button)
def intro_page():
    # Custom CSS for vertical and horizontal centering
    st.markdown("""
        <style>
            .centered-container {
                display: flex;
                justify-content: center;
                align-items: center;
                height: 90vh;  /* Vertically center the content */
                text-align: center;
            }
            .centered-content {
                display: flex;
                flex-direction: column;
                align-items: center;
                width: 100%;
                max-width: 400px;
                margin: 0 auto;
            }
            .centered-content h1 {
                margin-bottom: 20px;
            }
            .centered-content input {
                width: 100%;
                height: 40px;
                font-size: 20px;
                margin-bottom: 20px;
            }
        </style>
        <div class="centered-container">
            <div class="centered-content">
                <h1>Welcome to the Chat App</h1>
                <p>Please enter your ID to continue</p>
            </div>
        </div>
    """, unsafe_allow_html=True)

    # Input field and button using Streamlit's form
    with st.form("user_id_form", clear_on_submit=True):
        user_id = st.text_input("Enter your ID", "")
        submitted = st.form_submit_button("Continue")

    if submitted:
        if user_id:  # Only proceed if an ID is entered
            st.session_state.user_id = user_id
            go_to_chat()  # Navigate to the chat page
        else:
            st.warning("Please enter a valid ID.")

if 'query_input' not in st.session_state:
    st.session_state.query_input = ""


def chat_page():
    st.header("RAG File Assistant", divider="red")

    # File uploader
    uploaded_files = st.file_uploader("Please Upload file", type=["png", "jpg", "jpeg", "pdf", "txt", "docx", "doc", "ppt", "pptx", "xls", "xlsx", "csv"], accept_multiple_files=True)
    
    # Detect if files have been uploaded and update session state
    if uploaded_files and uploaded_files != st.session_state.files_uploaded:
        st.session_state.files_uploaded = uploaded_files
        st.session_state.history = []  # Clear chat history
        st.session_state.file_contents = []  # Clear stored file contents
        st.session_state.file_titles = {}  # Clear stored file titles
        
        # Assign files to variables and titles
        file_vars = assign_files_to_vars(uploaded_files)


    col1, col2 = st.columns([0.3,0.7])

    with col1:
        input_container = st.container()

    with col2:
        chat_container = st.container()

    with input_container:
        # Display the uploaded files
        if st.session_state.file_titles:
            st.write("### Uploaded Files")
            for file_key, file_title in st.session_state.file_titles.items():
                st.write(f"{file_key}: {file_title}")
        
        
    # Display chat history
    with chat_container:
        st.write("### Conversation History")
        display_history()

        # Textarea


        # Textarea within a form for better handling
        with st.form(key='chat_form', clear_on_submit=True):
            query = st.text_area(
                "Ask your AI Assistant",
                max_chars=10000, 
                height=200, 
                key="query_input"
            )
            submit_button = st.form_submit_button(label="Send", on_click=send_message)



def send_message():
    query = st.session_state.query_input.strip()  # Get and trim the input
    
    if query:
        #st.session_state.history.append({"user": query})  # Append user message to history
        st.write(f"Message sent: {query}")  # Replace with actual processing logic
        
        # Process files if any are uploaded
        if st.session_state.files_uploaded:
            process_files(st.session_state.files_uploaded)
        
        # Handle the user's query using the uploaded files content
        handle_user_query(query)
        
        st.session_state.query_input = ''  # Clear the textarea
    else:
        st.warning("Please enter a message before sending.")

####################
### MAIN HANDLER ###
####################

if 'query_input' not in st.session_state:
    st.session_state.query_input = ""

# Initialize session state for page navigation
if 'page' not in st.session_state:
    st.session_state.page = 'intro'  # Start on the intro page

# Function to switch to the chat page
def go_to_chat():
    st.session_state.page = 'chat'

# Page navigation logic
if st.session_state.page == 'intro':
    intro_page()  # Display the intro page
elif st.session_state.page == 'chat':
    chat_page()  # Display the chat page