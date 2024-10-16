import streamlit as st

# TODO: Deprecated, migrate to langchain_openai.ChatOpenAI
from langchain_community.chat_models import ChatOpenAI

from langchain.schema.messages import HumanMessage

import pandas as pd
from docx import Document

# Library for data extraction,analysis, conversion & manipulation of PDFs. 
import fitz  # PyMuPDF

# Optical character recognition library
import pytesseract

from PIL import Image
from pptx import Presentation
import csv
from zipfile import BadZipFile  # Import BadZipFile from zipfile module

from io import BytesIO


# Initialize LangChain models
chain_gpt_35 = ChatOpenAI(model="gpt-3.5-turbo")


# Sidebar for API key input
#with st.sidebar:
#    api_key = st.text_input("Enter your OpenAI API Key", type="password")
#    if api_key:
#        os.environ["OPENAI_API_KEY"] = api_key


# Assign files to variables and titles
def assign_files_to_vars(files):
    file_vars = {}
    for i, file in enumerate(files):
        file_key = f'file{i+1}'
        file_vars[file_key] = file
        st.session_state.file_titles[file_key] = file.name  # Store file title
    return file_vars



# Function to append response to chat history
def append_response_to_history(user_query, response):
    # Append both user query and assistant response to history
    st.session_state.history.append({"user": user_query, "assistant": response.content})

# Function to summarize content if it's too long
def summarize_content(content):
    # If the content length exceeds 2048 characters, summarize it
    if len(content) > 16384:
        response = chain_gpt_35.invoke([HumanMessage(content=f"Summarize the following content:\n\n{content[:16000]}")])
        return response.content
    return content

# Function to process files
def process_files(files):
    data_found = False
    file_vars = assign_files_to_vars(files)
    for file_key, file in file_vars.items():
        file_type = file.type
        if file_type in ['image/png', 'image/jpeg']:
            data_found |= img_reader(file, file_key)
        elif file_type == 'application/pdf':
            data_found |= pdf_reader(file, file_key)
        elif file_type == 'text/plain':
            data_found |= txt_reader(file, file_key)
        elif file_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
            data_found |= doc_reader(file, file_key)
        elif file_type in ['application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
            data_found |= excel_reader(file, file_key)
        elif file_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
            data_found |= ppt_reader(file, file_key)
        elif file_type == 'text/csv':
            data_found |= csv_reader(file, file_key)
    
    if not data_found:
        #st.warning("The provided files do not meet your requirements.")  # Use warning instead of adding to history
        pass

# File Readers with content extraction

# Image reader for image files containing text
def img_reader(file, file_key):
    img = Image.open(file)
    text = pytesseract.image_to_string(img)
    st.session_state.file_contents.append(f"{file_key}: {summarize_content(text)}")

    return True

# PDF reader using fitz library
def pdf_reader(file, file_key):
    try:
        file_stream = BytesIO(file.read())
        if file_stream.getbuffer().nbytes == 0:
            #st.error("The provided PDF file is empty.")  # Use error instead of adding to history
            return False
        
        pdf_document = fitz.open(stream=file_stream, filetype="pdf")
        full_text = ""
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            full_text += page.get_text()
        pdf_document.close()
        
        if full_text.strip() == "":
            st.error("The provided PDF file does not contain readable text.")  # Use error instead of adding to history
            return False
        
        st.session_state.file_contents.append(f"{file_key}: {summarize_content(full_text)}")
        return True

    except fitz.EmptyFileError:
        st.error("The provided PDF file is empty or corrupted.")  # Use error instead of adding to history
        return False

def txt_reader(file, file_key):
    text = file.read().decode("utf-8")
    if text.strip() == "":
        return False
    st.session_state.file_contents.append(f"{file_key}: {summarize_content(text)}")
    return True

def doc_reader(file, file_key):
    try:
        doc = Document(BytesIO(file.read()))
        text = "\n".join([para.text for para in doc.paragraphs])
        if text.strip() == "":
            #st.error("The provided Word document is empty or does not contain readable text.")  # Use error instead of adding to history
            return False
        st.session_state.file_contents.append(f"{file_key}: {summarize_content(text)}")
        return True

    except BadZipFile:
        st.error("The provided file could not be read as a valid Word document (BadZipFile).")  # Use error instead of adding to history
        return False

def excel_reader(file, file_key):
    df = pd.read_excel(file)
    text = df.to_string()
    if text.strip() == "":
        return False
    st.session_state.file_contents.append(f"{file_key}: {summarize_content(text)}")
    return True

def ppt_reader(file, file_key):
    presentation = Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    if text.strip() == "":
        return False
    st.session_state.file_contents.append(f"{file_key}: {summarize_content(text)}")
    return True

def csv_reader(file, file_key):
    df = pd.read_csv(file)
    text = df.to_string()
    if text.strip() == "":
        return False
    st.session_state.file_contents.append(f"{file_key}: {summarize_content(text)}")
    return True

# Function to handle user queries and invoke models
def handle_user_query(user_query):
    combined_content = "\n\n".join(st.session_state.file_contents)  # Combine all file contents
    summarized_content = summarize_content(combined_content)  # Summarize combined content
    prompt = f"Based on the following uploaded document contents:\n\n{summarized_content}\n\nUser's question: {user_query}"

    # Transform a single input into an output
    # HumanMessage represents a message with role "user"
    response = chain_gpt_35.invoke([HumanMessage(content=prompt)])


    append_response_to_history(user_query, response)


# Display chat history

#def display_history():
#    for message in st.session_state.history[-5:]:  # Display only the last 5 messages to manage token limit
#        if message['user']:  # Only display messages that are not empty or system messages
#            st.write(f"**User:** {message['user']}")
#            st.write(f"**Assistant:** {message['assistant']}")


def display_history():
    for message in st.session_state.history[-5:]:  # Display only the last 5 messages to manage token limit
        if message['user']:  # Only display messages that are not empty or system messages
            st.markdown(f"""
                <p style="margin-bottom:2px;font-weight:bold;">User:</p>
                <div style='background-color:#d1f0ff; padding:10px; border-radius:10px; margin-bottom:10px;'>
                    <p>{message['user']}</p>
                </div>
            """, unsafe_allow_html=True)
            
            st.markdown(f"""
                <p style="margin-bottom:2px;font-weight:bold;">Assistant:</p>
                <div style='background-color:#f7f7f7; padding:10px; border-radius:10px; margin-bottom:10px;'>
                    <p>{message['assistant']}</p>
                </div>
            """, unsafe_allow_html=True)

if 'query_input' not in st.session_state:
    st.session_state.query_input = ""


def send_message():
    query = st.session_state.query_input.strip()  # Get and trim the input
    
    if query:
        #st.session_state.history.append({"user": query})  # Append user message to history
        st.write(f"Message sent: {query}")  # Replace with actual processing logic
        
        # Process files if any are uploaded
        if st.session_state.files_uploaded:
            process_files(st.session_state.files_uploaded)
        
        # Handle the user's query using the uploaded files content
        handle_user_query(query)
        
        st.session_state.query_input = ''  # Clear the textarea
    else:
        st.warning("Please enter a message before sending.")